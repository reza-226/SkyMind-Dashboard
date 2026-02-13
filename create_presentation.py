from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import random

# ==================== رنگ‌های دقیق ====================
C_BG_DARK   = RGBColor(0x0a, 0x16, 0x28)
C_BG_MED    = RGBColor(0x1a, 0x29, 0x42)
C_CYAN      = RGBColor(0x00, 0xd9, 0xff)
C_ORANGE    = RGBColor(0xff, 0x6b, 0x35)
C_WHITE     = RGBColor(0xff, 0xff, 0xff)
C_GRAY      = RGBColor(0xb4, 0xb4, 0xb4)

# ==================== گرادیانت ====================
def add_gradient_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_angle = 135.0
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[0].color.rgb = C_BG_DARK
    fill.gradient_stops[1].position = 1.0
    fill.gradient_stops[1].color.rgb = C_BG_MED

# ==================== شبکه دیجیتال HUD ====================
def add_network_hud(slide):
    """شبکه چندضلعی با نقاط و خطوط نورانی"""
    random.seed(42)
    sh = slide.shapes
    
    # ایجاد نقاط شبکه
    nodes = []
    for i in range(12):
        x = Inches(0.3 + i * 0.85)
        y = Inches(0.2 + random.uniform(-0.12, 0.12))
        nodes.append((x, y))
        
        # نقطه نورانی
        dot = sh.add_shape(9, x, y, Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_CYAN
        dot.line.color.rgb = C_CYAN
        dot.line.width = Pt(1)
    
    # خطوط اتصال افقی
    for i in range(len(nodes) - 1):
        x1, y1 = nodes[i][0] + Inches(0.04), nodes[i][1] + Inches(0.04)
        x2, y2 = nodes[i+1][0] + Inches(0.04), nodes[i+1][1] + Inches(0.04)
        line = sh.add_connector(1, x1, y1, x2, y2)
        line.line.color.rgb = C_CYAN
        line.line.width = Pt(0.8)
    
    # خطوط مورب (چندضلعی)
    for i in [0, 2, 5, 7, 9]:
        if i + 2 < len(nodes):
            x1, y1 = nodes[i][0] + Inches(0.04), nodes[i][1] + Inches(0.04)
            x2, y2 = nodes[i+2][0] + Inches(0.04), nodes[i+2][1] + Inches(0.04)
            line = sh.add_connector(1, x1, y1, x2, y2)
            line.line.color.rgb = C_CYAN
            line.line.width = Pt(0.5)
    
    # متن‌های فنی کوچک (HUD)
    for i in [1, 4, 7, 10]:
        x = nodes[i][0] - Inches(0.15)
        y = nodes[i][1] - Inches(0.25)
        tb = sh.add_textbox(x, y, Inches(0.4), Inches(0.15))
        tf = tb.text_frame
        tf.text = f"NODE-{i:02d}"
        p = tf.paragraphs[0]
        p.font.size = Pt(6)
        p.font.color.rgb = C_CYAN
        p.font.name = "Consolas"
        p.alignment = PP_ALIGN.CENTER

# ==================== خطوط توپوگرافی ====================
def add_topo_curves(slide):
    """خطوط کانتور پایین"""
    sh = slide.shapes
    for i in range(6):
        y = Inches(6.4 + i * 0.15)
        line = sh.add_connector(1, Inches(0), y, Inches(10), y)
        line.line.color.rgb = C_CYAN
        line.line.width = Pt(0.7)

# ==================== پهپاد Wireframe ====================
def add_drone_wireframe(slide):
    """پهپاد مرکزی با طراحی مهندسی"""
    sh = slide.shapes
    
    # بدنه اصلی (مستطیل مرکزی)
    body = sh.add_shape(1, Inches(4.5), Inches(3), Inches(1), Inches(0.6))
    body.fill.background()
    body.line.color.rgb = C_WHITE
    body.line.width = Pt(2)
    
    # بازوهای پهپاد (4 بازو)
    arm_positions = [
        (Inches(4.2), Inches(2.6)),  # بالا-چپ
        (Inches(5.3), Inches(2.6)),  # بالا-راست
        (Inches(4.2), Inches(3.4)),  # پایین-چپ
        (Inches(5.3), Inches(3.4)),  # پایین-راست
    ]
    
    for x, y in arm_positions:
        arm = sh.add_shape(1, x, y, Inches(0.15), Inches(0.15))
        arm.fill.solid()
        arm.fill.fore_color.rgb = C_CYAN
        arm.line.color.rgb = C_CYAN
        arm.line.width = Pt(1.5)
    
    # دوربین/حسگر جلو
    cam = sh.add_shape(9, Inches(4.8), Inches(2.85), Inches(0.2), Inches(0.15))
    cam.fill.solid()
    cam.fill.fore_color.rgb = C_ORANGE
    cam.line.color.rgb = C_ORANGE
    cam.line.width = Pt(1)
    
    # متن‌های فنی اطراف پهپاد
    labels = [
        (Inches(4.1), Inches(2.4), "ROTOR-1"),
        (Inches(5.4), Inches(2.4), "ROTOR-2"),
        (Inches(4.1), Inches(3.7), "ROTOR-3"),
        (Inches(5.4), Inches(3.7), "ROTOR-4"),
    ]
    
    for x, y, text in labels:
        tb = sh.add_textbox(x, y, Inches(0.5), Inches(0.15))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(7)
        p.font.color.rgb = C_CYAN
        p.font.name = "Consolas"

# ==================== ابر داده ====================
def add_data_cloud(slide):
    """ابر چندضلعی با ستون‌های نور"""
    sh = slide.shapes
    
    # شکل ابر (چندضلعی تقریبی)
    cloud = sh.add_shape(17, Inches(3.8), Inches(1.2), Inches(2.4), Inches(1.2))
    cloud.fill.solid()
    cloud.fill.fore_color.rgb = RGBColor(0x00, 0x4d, 0x66)
    cloud.line.color.rgb = C_CYAN
    cloud.line.width = Pt(2)
    
    ct = cloud.text_frame
    ct.text = "DATA CLOUD\n01010110"
    ct.paragraphs[0].font.size = Pt(11)
    ct.paragraphs[0].font.bold = True
    ct.paragraphs[0].font.color.rgb = C_WHITE
    ct.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    ct.paragraphs[1].font.size = Pt(8)
    ct.paragraphs[1].font.name = "Courier New"
    ct.paragraphs[1].font.color.rgb = C_CYAN
    ct.paragraphs[1].alignment = PP_ALIGN.CENTER
    
    # ستون‌های نور عمودی از ابر به پهپاد
    for dx in [-0.3, 0, 0.3]:
        x = Inches(5 + dx)
        line = sh.add_connector(1, x, Inches(2.4), x, Inches(3))
        line.line.color.rgb = C_CYAN
        line.line.width = Pt(1.5)
        line.line.dash_style = 3  # نقطه‌چین

# ==================== شهر 3D Wireframe ====================
def add_city_wireframe(slide):
    """ساختمان‌های شهری با خطوط نورانی"""
    sh = slide.shapes
    
    # ساختمان‌های مختلف (ارتفاع متفاوت)
    buildings = [
        (Inches(0.5), Inches(5.2), Inches(0.6), Inches(1.8)),
        (Inches(1.3), Inches(4.8), Inches(0.5), Inches(2.2)),
        (Inches(2.0), Inches(5.5), Inches(0.4), Inches(1.5)),
        (Inches(7.5), Inches(5.0), Inches(0.7), Inches(2)),
        (Inches(8.4), Inches(5.3), Inches(0.5), Inches(1.7)),
        (Inches(9.2), Inches(5.6), Inches(0.4), Inches(1.4)),
    ]
    
    for x, y, w, h in buildings:
        bldg = sh.add_shape(1, x, y, w, h)
        bldg.fill.background()
        bldg.line.color.rgb = C_WHITE
        bldg.line.width = Pt(1.5)
        
        # نقطه نورانی روی ساختمان
        dot_x = x + w / 2 - Inches(0.04)
        dot_y = y - Inches(0.1)
        dot = sh.add_shape(9, dot_x, dot_y, Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_CYAN
        dot.line.color.rgb = C_CYAN
    
    # خطوط اتصال شهری
    connections = [
        ((Inches(0.8), Inches(5.2)), (Inches(1.55), Inches(4.8))),
        ((Inches(1.55), Inches(4.8)), (Inches(2.2), Inches(5.5))),
        ((Inches(7.85), Inches(5)), (Inches(8.65), Inches(5.3))),
        ((Inches(8.65), Inches(5.3)), (Inches(9.4), Inches(5.6))),
    ]
    
    for (x1, y1), (x2, y2) in connections:
        line = sh.add_connector(1, x1, y1, x2, y2)
        line.line.color.rgb = C_CYAN
        line.line.width = Pt(1)

# ==================== خطوط داده پهپاد-شهر ====================
def add_data_streams(slide):
    """خطوط منحنی از پهپاد به شهر"""
    sh = slide.shapes
    
    streams = [
        ((Inches(4.7), Inches(3.6)), (Inches(1.5), Inches(5.5))),
        ((Inches(5.3), Inches(3.6)), (Inches(8.5), Inches(5.8))),
    ]
    
    for (x1, y1), (x2, y2) in streams:
        line = sh.add_connector(2, x1, y1, x2, y2)  # 2 = curved
        line.line.color.rgb = C_CYAN
        line.line.width = Pt(1.2)

# ==================== اسلاید عنوان ====================
def slide_title_main(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)
    add_network_hud(slide)
    add_topo_curves(slide)
    add_drone_wireframe(slide)
    add_data_cloud(slide)
    add_city_wireframe(slide)
    add_data_streams(slide)
    
    sh = slide.shapes
    
    # عنوان اصلی
    tb = sh.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "SkyMind: بهینه‌سازی تخلیه محاسباتی"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "در شبکه‌های پهپادی"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.alignment = PP_ALIGN.CENTER
    
    # زیرعنوان
    tb2 = sh.add_textbox(Inches(1), Inches(2.1), Inches(8), Inches(0.6))
    p3 = tb2.text_frame.paragraphs[0]
    p3.text = "رویکردی نوین با استفاده از یادگیری تقویتی چندعامله (MADDPG)"
    p3.font.size = Pt(20)
    p3.font.color.rgb = C_CYAN
    p3.alignment = PP_ALIGN.CENTER
    
    # پاورقی
    tb3 = sh.add_textbox(Inches(2), Inches(6.9), Inches(6), Inches(0.4))
    p4 = tb3.text_frame.paragraphs[0]
    p4.text = "دفاع پایان‌نامه کارشناسی ارشد | [نام دانشجو]"
    p4.font.size = Pt(14)
    p4.font.color.rgb = C_GRAY
    p4.alignment = PP_ALIGN.CENTER

# ==================== اسلاید محتوایی ====================
def slide_content_example(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)
    add_network_hud(slide)
    add_topo_curves(slide)
    
    sh = slide.shapes
    
    # عنوان
    tb = sh.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = "انگیزه تحقیق: شکاف میان داده و منابع"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # خط نارنجی
    line = sh.add_shape(1, Inches(2), Inches(1.05), Inches(6), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = C_ORANGE
    line.line.fill.background()
    
    # محتوا
    content = """
🔹 چالش اصلی:
   • رشد انفجاری داده‌های IoT (44 → 175 زتابایت)
   • نیاز به پاسخ آنی (10–50 میلی‌ثانیه)
   • محدودیت منابع پهپادها (باتری، پردازش)

🔹 راه‌حل پیشنهادی:
   استفاده از پهپادها به عنوان سرورهای لبه متحرک
   با بهینه‌سازی هوشمند توسط MADDPG
"""
    
    tb2 = sh.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(4))
    tf = tb2.text_frame
    tf.word_wrap = True
    tf.text = content.strip()
    
    for para in tf.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = C_WHITE
        para.space_after = Pt(12)

# ==================== ساخت فایل ====================
def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    print("➤ اسلاید 1: صفحه عنوان (طراحی PDF اصلی)...")
    slide_title_main(prs)
    
    print("➤ اسلاید 2: نمونه محتوایی...")
    slide_content_example(prs)
    
    out = "SkyMind_UAV_Presentation.pptx"
    prs.save(out)
    print(f"\n✅ فایل ذخیره شد: {out}")

if __name__ == "__main__":
    build_presentation()
